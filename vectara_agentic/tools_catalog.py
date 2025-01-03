"""
This module contains the tools catalog for the Vectara Agentic.
"""
from typing import List, Dict, Any
from functools import lru_cache
from pydantic import Field
from pptx import Presentation as CreatePresentation # This is the function to create presentations
from pptx.presentation import Presentation # This is the actual class type
from io import BytesIO
import requests

from .types import LLMRole
from .utils import get_llm

req_session = requests.Session()

get_headers = {
    "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64; rv:98.0) Gecko/20100101 Firefox/98.0",
    "Accept": "text/html,application/xhtml+xml,application/xml;q=0.9,image/avif,image/webp,*/*;q=0.8",
    "Accept-Language": "en-US,en;q=0.5",
    "Accept-Encoding": "gzip, deflate",
    "Connection": "keep-alive",
}


#
# Standard Tools
#
@lru_cache(maxsize=None)
def summarize_text(
    text: str = Field(description="the original text."),
    expertise: str = Field(
        description="the expertise to apply to the summarization.",
    ),
) -> str:
    """
    This is a helper tool.
    Use this tool to summarize text using a given expertise
    with no more than summary_max_length characters.

    Args:
        text (str): The original text.
        expertise (str): The expertise to apply to the summarization.

    Returns:
        str: The summarized text.
    """
    if not isinstance(expertise, str):
        return "Please provide a valid string for expertise."
    if not isinstance(text, str):
        return "Please provide a valid string for text."
    expertise = "general" if len(expertise) < 3 else expertise.lower()
    prompt = f"As an expert in {expertise}, summarize the provided text"
    prompt += " into a concise summary."
    prompt += f"\noriginal text: {text}\nsummary:"
    llm = get_llm(LLMRole.TOOL)
    response = llm.complete(prompt)
    return response.text


@lru_cache(maxsize=None)
def rephrase_text(
    text: str = Field(description="the original text."),
    instructions: str = Field(description="the specific instructions for how to rephrase the text."),
) -> str:
    """
    This is a helper tool.
    Use this tool to rephrase the text according to the provided instructions.
    For example, instructions could be "as a 5 year old would say it."

    Args:
        text (str): The original text.
        instructions (str): The specific instructions for how to rephrase the text.

    Returns:
        str: The rephrased text.
    """
    prompt = f"""
    Rephrase the provided text according to the following instructions: {instructions}.
    If the input is Markdown, keep the output in Markdown as well.
    original text: {text}
    rephrased text:
    """
    llm = get_llm(LLMRole.TOOL)
    response = llm.complete(prompt)
    return response.text


@lru_cache(maxsize=None)
def critique_text(
    text: str = Field(description="the original text."),
    role: str = Field(default=None, description="the role of the person providing critique."),
    point_of_view: str = Field(default=None, description="the point of view with which to provide critique."),
) -> str:
    """
    This is a helper tool.
    Critique the text from the specified point of view.

    Args:
        text (str): The original text.
        role (str): The role of the person providing critique.
        point_of_view (str): The point of view with which to provide critique.

    Returns:
        str: The critique of the text.
    """
    if role:
        prompt = f"As a {role}, critique the provided text from the point of view of {point_of_view}."
    else:
        prompt = f"Critique the provided text from the point of view of {point_of_view}."
    prompt += "Structure the critique as bullet points.\n"
    prompt += f"Original text: {text}\nCritique:"
    llm = get_llm(LLMRole.TOOL)
    response = llm.complete(prompt)
    return response.text


#
# Guardrails tool: returns list of topics to avoid
#
def get_bad_topics() -> List[str]:
    """
    Get the list of topics to avoid in the response.
    """
    return [
        "politics",
        "religion",
        "violence",
        "hate speech",
        "adult content",
        "illegal activities",
    ]

#
# Document Template tools
#
def get_presentation_from_path(
    file_path: str
    ) -> Presentation:
    """"
    Reads in the file path of a PowerPoint presentation template and returns a Presentation object that can be filled in using other tools.

    Args:
        file_path (str): The file path of the presentation to edit.

    Returns:
        Presentation: the editable presentation object
    """
    return CreatePresentation(file_path)

def get_presentation_field(
    presentation: Presentation,
    slide_idx: int = 0
    ) -> Dict:
    """
    Returns information about the next field to fill out in a presentation template.

    Args:
        presentation (Presentation): The presentation to get the field information from.
        slide_idx (int): The slide index of the previous field. Defaults to 1.

    Returns:
        Dict: Information about the field, including the question to answer to populate the field,
                the context in the presentation where the answer to the question will be inserted, and
                other variables specifying the location of the field in the presentation.
    """
    for slide_i, slide in enumerate(presentation.slides):
        if slide_i < slide_idx:
            continue

        for element_i, element in enumerate(slide.shapes):
            if element.has_text_frame:
                for paragraph_i, paragraph in enumerate(element.text_frame.paragraphs):
                    line = paragraph.text
                    if (line.find("{{") != -1) and (line.find("}}") != -1):
                        return {
                            "question": line[line.find("{{")+2:line.find("}}")],
                            "context": line,
                            "slide_idx": slide_i,
                            "element_idx": element_i,
                            "paragraph_idx": paragraph_i
                        }
    return {
        "question": "No more questions to answer in presentation."
    }

def fill_presentation_field(
    presentation: Presentation,
    question: str,
    answer: str,
    slide_idx: int,
    element_idx: int,
    paragraph_idx: int
    ) -> Presentation:
    """
    Fills in the specified field location with the answer to the question.

    Args:
        presentation (Presentation): The presentation to edit.
        question (str): The question in the field that should be replaced.
        answer (str): The answer to the question that will be populated in the field.
        slide_idx (int): The index of the slide where the field is located.
        element_idx (int): The index of the element where the field is located.
        paragraph_idx (int): The index of the paragraph where the field is located.

    Returns:
        Presentation: The presentation with the specified field populated.
    """

    try:
        paragraph = presentation.slides[slide_idx].shapes[element_idx].text_frame.paragraphs[paragraph_idx]

        if 'NEEDS REVIEW' in answer:
            paragraph.text = paragraph.text.replace(f"{{{question}}}", answer)
        else:
            paragraph.text = paragraph.text.replace(f"{{{{{question}}}}}", answer)

        return presentation

    except Exception as e:
        return f"Error: {str(e)} occurred while trying to populate presentation field."

def get_output_presentation_file(
    presentation: Presentation
    ) -> BytesIO:
    """
    Returns a BytesIO object that can be downloaded by a user from a filled presentation.

    Args:
        presentation (Presentation): The presentation object.

    Returns:
        BytesIO: The downloadable file.
    """
    output = BytesIO()
    presentation.save(output)
    output.seek(0)

    return output


# def populate_ppt(
#         file_path: str,
#         topic: str
#     ) -> BytesIO:
#     """
#     Fills in a PowerPoint presentation template by asking a series of questions and populating the template based on retrieved answers.

#     Args:
#         file_path (str): The file path of the presentation to edit.
#         topic (str): A short description about the topic of the presentation.
#     """
#     ppt = Presentation(file_path)

#     # Create the prompt for the LLM
#     prompt_template = "Your job is to answer questions that will be used to populate a presentation template about {topic}."
#     prompt_template += f"\nYou have access to the following tools: {tool_info}"
#     prompt_template += "\n Use these tools to answer the question {question} in a concise manner that will go well in the following paragraph:\n{context}."
#     prompt_template += "\nIf you do not have enough information to answer the question, respond with 'NEEDS REVIEW'."
#     llm = get_llm(LLMRole.TOOL) # This is a problem because the tool LLM does not have access to the other tools.
    
#     for slide in ppt.slides:
#         for element in slide.shapes:
#             if element.has_text_frame:
#                 for paragraph in element.text_frame.paragraphs:
#                     line = paragraph.text

#                     # While there are still questions in this paragraph:
#                     while (line.find("{{") != -1) and (line.find("}}") != -1):
#                         print(f"DEBUG: CURRENT PARAGRAPH IS {line}")
#                         # Get the question from the paragraph
#                         question = line[line.find("{{")+2:line.find("}}")]

#                         print(f"DEBUG: ANSWERING QUESTION {question}")

#                         # Call the main llm to try to answer the question (create a prompt with the question and the context, give instructions for default answer)
#                         prompt = prompt_template.format(topic=topic, question=question, context=line)
#                         response = llm.complete(prompt)
#                         answer = response.text

#                         print(f"DEBUG: RECEIVED ANSWER {answer}")

#                         # Replace the question with the answer to the question.
#                         # If the "I don't know"-based response is provided, format the text (with highlighting) so that the user knows to review the information themselves.
#                         if 'NEEDS REVIEW' in answer:
#                             paragraph.text = line.replace(f"{{{question}}}", answer) # Leave in one set of {}
#                         else:
#                             paragraph.text = line.replace(f"{{{{{question}}}}}", answer)
#                         line = paragraph.text

#     output = BytesIO()
#     ppt.save(output)
#     output.seek(0)

#     return output
